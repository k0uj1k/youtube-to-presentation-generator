import os
import re
import shutil
import subprocess
import time
import uuid
import cv2
import numpy as np
import functools
from ja_sentence_segmenter.common.pipeline import make_pipeline
from ja_sentence_segmenter.concatenate.simple_concatenator import concatenate_matching
from ja_sentence_segmenter.normalize.neologd_normalizer import normalize
from ja_sentence_segmenter.split.simple_splitter import split_newline, split_punctuation
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api._errors import TranscriptsDisabled, NoTranscriptFound
import yt_dlp
import budoux
from pptx import Presentation
from pptx.util import Inches, Pt
# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor
import unicodedata
from dotenv import load_dotenv
from .gemini_service import GeminiSummarizer
from .task_manager import TaskCancelledException

# .env ファイルを読み込む
load_dotenv()

# 一時ファイルを保存するディレクトリ
TEMP_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "temp_data")
os.makedirs(TEMP_DIR, exist_ok=True)

# 変化レベル（1〜10）を MAD 閾値に変換するテーブル
# レベル 1 = 非常に敏感（わずかな変化も検知）、10 = 鈍感（大きな変化のみ検知）
_CHANGE_LEVEL_TO_THRESHOLD = {
    1:   1.0,
    2:   3.0,
    3:   6.0,
    4:  10.0,
    5:  15.0,
    6:  25.0,
    7:  40.0,
    8:  60.0,
    9:  80.0,
   10: 100.0,
}


def split_japanese_sentences(text: str) -> str:
    """日本語のテキストを ja_sentence_segmenter を用いて文ごとに分割し、
    改行で区切られた読みやすい文字列に整形する。
    句読点がない自動生成字幕などの場合は、スペースで分割するフォールバックを行う。
    """
    if not text:
        return text
        
    # 句読点で区切るルール
    split_punc = functools.partial(split_punctuation, punctuations=r"。!?")
    
    # 連結ルール（「の」で終わる文などを結合）
    concat_tail_no = functools.partial(
        concatenate_matching, 
        former_matching_rule=r"^(?P<result>.+)(の)$", 
        remove_former_matched=False
    )
    
    # パイプライン構築
    segmenter = make_pipeline(normalize, split_newline, concat_tail_no, split_punc)
    
    # セグメント実行
    sentences = list(segmenter(text))
    
    # もし文が分割されなかった（句読点がなく、1つの長い文になっている）場合
    # 自動字幕を想定し、スペース（半角/全角）での分割を試みる
    if len(sentences) <= 1:
        spaced_parts = re.split(r'\s+', text)
        spaced_parts = [p.strip() for p in spaced_parts if p.strip()]
        if len(spaced_parts) > 1:
            return "\n".join(spaced_parts)
            
    return "\n".join(sentences).strip()


def apply_budoux_layout(text: str, max_line_len: int = 25) -> str:
    """BudouX を使用して、日本語テキストが不自然な位置で改行されないように、
    適切なフレーズ境界で改行（\n）を挿入して整形する。
    """
    if not text:
        return text
        
    parser = budoux.load_default_japanese_parser()
    lines = text.splitlines()
    formatted_lines = []
    
    for line in lines:
        if len(line) <= max_line_len:
            formatted_lines.append(line)
            continue
            
        # フレーズに分割
        phrases = parser.parse(line)
        
        current_line = ""
        line_parts = []
        
        for phrase in phrases:
            # 現在の行に追加すると上限文字数を超える場合、改行
            if len(current_line) + len(phrase) > max_line_len:
                if current_line:
                    line_parts.append(current_line)
                current_line = phrase
            else:
                current_line += phrase
                
        if current_line:
            line_parts.append(current_line)
            
        formatted_lines.extend(line_parts)
        
    return "\n".join(formatted_lines)


def is_japanese(text: str) -> bool:
    """テキストに日本語の文字（ひらがな、カタカナ、漢字）が含まれているか判定する。"""
    if not text:
        return False
    # ひらがな、カタカナ、CJK統合漢字の範囲を正規表現でチェック
    return bool(re.search(r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF]', text))


def clean_english_text(text: str) -> str:
    """英文テキストをクレンジングし、改行やハイフン、余分なスペースを整形する。"""
    if not text:
        return text

    # 1. 改行で分断された単語のハイフン(-)を連結
    text = re.sub(r'-\s*\n', '', text)

    # 2. 連続する改行やスペースを単一のスペースに置換
    # 2つ以上の改行を段落の区切りとして保持しつつ、単一の改行をスペースに変換
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)

    # 3. 余分なスペースの削除
    # 改行以外の連続する空白文字を1つのスペースにする
    text = re.sub(r'[^\S\n]+', ' ', text)
    # 各行の先頭と末尾の空白をトリム
    lines = [line.strip() for line in text.split('\n')]
    # 連続する空行を整理する
    cleaned_text = '\n'.join(lines)
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)

    return cleaned_text.strip()


def format_slide_text(text: str) -> str:
    """テキスト言語（日本語/その他）を判定し、最適な整形処理を適用する。"""
    if not text:
        return text

    if is_japanese(text):
        # 日本語の場合は文分割 + BudouX 禁則処理
        text = split_japanese_sentences(text)
        text = apply_budoux_layout(text)
    else:
        # 英語など日本語以外の場合は英文クレンジング
        text = clean_english_text(text)

    return text


def sanitize_filename(title: str, max_length: int = 200) -> str:
    """
    動画タイトルをファイルシステムセーフなファイル名に変換する。
    
    - Unicode正規化（NFKC）
    - ファイルシステムで禁止されている文字を削除
    - 先頭/末尾の空白やドットを削除
    - 最大文字数を制限
    """
    if not title:
        return "presentation"
    
    # Unicode正規化（濁音などを統一）
    title = unicodedata.normalize('NFKC', title)
    
    # ファイルシステムで禁止されている文字を削除・置換
    # Windows: < > : " / \ | ? *
    # macOS/Linux: /
    forbidden_chars = r'[<>:"/\\|?*]'
    sanitized = re.sub(forbidden_chars, '', title)
    
    # 連続する空白をシングルスペースに
    sanitized = re.sub(r'\s+', ' ', sanitized)
    
    # 先頭/末尾の空白やドットを削除（ファイルシステムの要件）
    sanitized = sanitized.strip('. ')
    
    # 最大文字数を制限（.pptx を考慮）
    if len(sanitized) > max_length:
        sanitized = sanitized[:max_length].strip()
    
    # 万が一空文字列になった場合のフォールバック
    if not sanitized:
        sanitized = "presentation"
    
    return sanitized


def format_timestamp(seconds: float) -> str:
    """秒数を mm:ss 形式の文字列に変換する。"""
    minutes = int(seconds // 60)
    secs = int(seconds % 60)
    return f"{minutes:02d}:{secs:02d}"


def get_scene_text(transcript: list, current_time: float, next_time: float) -> str:
    """指定区間に含まれる字幕を収集し、スライド向けに整形して返す。"""
    if not transcript:
        return ""

    slide_text_list = []
    for entry in transcript:
        start = entry["start"]
        if current_time <= start < next_time:
            slide_text_list.append(entry["text"])

    return format_slide_text("\n".join(slide_text_list).strip())


def create_markdown_package(
    title: str,
    scenes: list,
    transcript: list,
    task_temp_dir: str,
    safe_title: str,
    url: str | None = None,
    ai_summary_enabled: bool = False,
    task_state=None
) -> tuple[str, str, list[str]]:
    """Markdown ファイルと画像フォルダを生成する。"""
    asset_dirname = "images"
    assets_dir = os.path.join(task_temp_dir, asset_dirname)
    markdown_filename = f"{safe_title}.md"
    markdown_path = os.path.join(task_temp_dir, markdown_filename)
    os.makedirs(assets_dir, exist_ok=True)
    asset_filenames = []

    summarizer = None
    if ai_summary_enabled:
        try:
            summarizer = GeminiSummarizer()
            print("✓ Gemini API が有効です。Markdown 要約を生成します。")
        except Exception as e:
            ai_summary_enabled = False
            print(f"ℹ Gemini 要約を無効化しました。({e})")

    lines = [
        f"# {title}",
        "",
        "YouTube動画から自動生成されたMarkdown資料",
        "",
        f"- スライド枚数: {len(scenes)}",
    ]
    if url:
        lines.extend([
            f"- URL: {url}",
        ])
    lines.append("")

    for i, scene in enumerate(scenes):
        if task_state:
            task_state.check_cancelled()
            progress = 80 + int((i + 1) / len(scenes) * 15)
            task_state.log(f"Markdown資料を作成中... ({i+1}/{len(scenes)}枚目)", progress)

        current_time = scene["timestamp"]
        next_time = scenes[i + 1]["timestamp"] if i + 1 < len(scenes) else float("inf")
        image_name = os.path.basename(scene["image_path"])
        packaged_image_path = os.path.join(assets_dir, image_name)
        shutil.copy2(scene["image_path"], packaged_image_path)
        asset_filenames.append(image_name)
        image_rel_path = f"./{asset_dirname}/{image_name}"

        slide_text = get_scene_text(transcript, current_time, next_time)
        display_text = slide_text or "(字幕なし)"

        lines.extend([
            f"## Slide {i + 1}",
            "",
            f"- 開始時刻: {format_timestamp(current_time)}",
            "",
            f"![Slide {i + 1}]({image_rel_path})",
            "",
        ])

        if ai_summary_enabled and summarizer and slide_text:
            try:
                summary_result = summarizer.summarize_slide_content(slide_text)
                key_points = summary_result.get("key_points") or []
                main_message = summary_result.get("main_message")

                if key_points:
                    lines.extend([
                        "### 要点",
                        "",
                    ])
                    lines.extend([f"- {format_slide_text(point)}" for point in key_points])
                    lines.append("")

                if main_message:
                    lines.extend([
                        "### 最も言いたいこと",
                        "",
                        main_message,
                        "",
                    ])
            except Exception as e:
                print(f"スライド {i+1} の要約生成に失敗: {e}")

        lines.extend([
            "### 文字起こし",
            "",
            display_text.replace("\n", "  \n"),
            "",
        ])

    with open(markdown_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"Markdown ファイルを保存しました: {markdown_path}")
    print(f"Markdown 画像フォルダを保存しました: {assets_dir}")
    return markdown_filename, asset_dirname, asset_filenames


def extract_video_id(url: str) -> str:
    """
    YouTubeのURLから動画IDを抽出する。
    """
    # 許可するホスト名の確認 (スキームが含まれている、またはプロトコル相対の場合)
    if "://" in url or url.startswith("//"):
        from urllib.parse import urlparse
        try:
            url_to_parse = "http:" + url if url.startswith("//") else url
            parsed_url = urlparse(url_to_parse)
            
            if parsed_url.scheme not in ["http", "https"]:
                raise ValueError("許可されていない URL スキームです。")
                
            netloc = parsed_url.netloc.lower()
            allowed_hosts = [
                "youtube.com", "www.youtube.com", "m.youtube.com",
                "youtu.be", "www.youtu.be", "music.youtube.com"
            ]
            
            is_valid_host = False
            for host in allowed_hosts:
                if netloc == host or netloc.endswith("." + host):
                    is_valid_host = True
                    break
                    
            if not is_valid_host:
                raise ValueError("YouTube 以外のホスト名からのダウンロードは許可されていません。")
        except Exception as e:
            if isinstance(e, ValueError):
                raise e
            raise ValueError("無効な URL 形式です。")

    # クエリパラメータ v= から抽出
    if "v=" in url:
        match = re.search(r"v=([0-9A-Za-z_-]{11})", url)
        if match:
            return match.group(1)

    # shortsのURLから抽出 (/shorts/xxx)
    if "/shorts/" in url:
        match = re.search(r"/shorts/([0-9A-Za-z_-]{11})", url)
        if match:
            return match.group(1)

    # youtu.be/xxx から抽出
    if "youtu.be/" in url:
        match = re.search(r"youtu\.be/([0-9A-Za-z_-]{11})", url)
        if match:
            return match.group(1)

    # その他のURLパターンからのフォールバック抽出
    patterns = [
        r"embed\/([0-9A-Za-z_-]{11})",
        r"\/([0-9A-Za-z_-]{11})",
        r"^([0-9A-Za-z_-]{11})$"  # 11文字の動画IDそのものを許可するパターン
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)

    raise ValueError("無効なYouTube URLまたは動画IDです。")


def get_transcript(url: str, video_id: str) -> list:
    """
    動画の字幕（文字起こし）を取得する。SubtitleFetcher のみで取得する。取得できない場合は空リストを返す。
    """
    from .subtitle_service import SubtitleFetcher
    fetcher = SubtitleFetcher(url, video_id, lang="ja")
    try:
        return fetcher.fetch()
    except Exception as e:
        print(f"字幕の取得に失敗しました: {e}")
        return []


def download_video(url: str, output_path: str) -> str:
    """
    yt-dlpを使用して画像抽出用に高精細な動画（例: 1080p）をダウンロードする。
    動画が既に存在する場合はダウンロードをスキップし、タイトルのみ取得して返す。
    """
    base_opts = {
        'quiet': True,
        'no_warnings': True,
    }
    
    try:
        with yt_dlp.YoutubeDL(base_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            title = info.get('title', 'YouTube Video')
    except Exception as e:
        error_msg = str(e)
        if "Video unavailable" in error_msg:
            raise ValueError("指定された動画は存在しないか、非公開、あるいは地域制限によりアクセスできません。")
        raise ValueError(f"動画情報の取得に失敗しました: {error_msg}")

    if not os.path.exists(output_path):
        print(f"動画をダウンロードします: {output_path}")
        ydl_opts = {
            'format': 'bestvideo[height<=1080][vcodec^=avc1]+bestaudio[ext=m4a]/best[height<=1080][vcodec^=avc1]/best[ext=mp4]',
            'outtmpl': output_path,
            'quiet': True,
            'no_warnings': True,
        }
        
        max_retries = 3
        for attempt in range(max_retries):
            try:
                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    ydl.extract_info(url, download=True)
                break
            except yt_dlp.utils.DownloadError as e:
                error_msg = str(e)
                if "Video unavailable" in error_msg:
                    raise ValueError("指定された動画は存在しないか、非公開、あるいは地域制限によりアクセスできません。")
                
                if attempt < max_retries - 1:
                    print(f"動画のダウンロードに失敗しました。3秒後にリトライします ({attempt + 1}/{max_retries} 回目の試行): {error_msg}")
                    time.sleep(3)
                else:
                    raise ValueError(f"動画のダウンロードに失敗しました: {error_msg}")
            except Exception as e:
                if attempt < max_retries - 1:
                    print(f"動画ダウンロード中に予期せぬエラーが発生しました。3秒後にリトライします ({attempt + 1}/{max_retries} 回目の試行): {e}")
                    time.sleep(3)
                else:
                    raise ValueError(f"動画のダウンロード処理中に予期せぬエラーが発生しました: {str(e)}")
    else:
        print(f"[INFO] 動画ファイルは既に存在します。ダウンロードをスキップします: {output_path}")

    return title


def get_keyframe_timestamps(video_path: str) -> list:
    """
    ffprobe を使用して Closed GOP の I フレーム（キーフレーム）の
    タイムスタンプ（秒）リストを返す。
    パケットレベルで K フラグ（key_frame）を持つフレームのみを抽出する。
    """
    cmd = [
        "ffprobe",
        "-v", "error",
        "-select_streams", "v:0",
        "-show_entries", "packet=pts_time,flags",
        "-of", "csv=print_section=0",
        video_path
    ]
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120
        )
    except FileNotFoundError:
        # ffprobe が見つからない場合は空リストを返す（フォールバックで対処）
        print("ffprobe が見つかりません。PATH に ffmpeg/ffprobe を追加してください。")
        return []
    except subprocess.TimeoutExpired:
        print("ffprobe がタイムアウトしました。")
        return []

    timestamps = []
    for line in result.stdout.splitlines():
        parts = line.strip().split(",")
        if len(parts) < 2:
            continue
        pts_str = parts[0]
        flags = parts[1] if len(parts) > 1 else ""
        # "K" フラグ = キーフレーム（I フレーム）
        if "K" in flags.upper():
            try:
                timestamps.append(float(pts_str))
            except ValueError:
                pass

    return sorted(set(timestamps))


def get_keyframe_timestamps_cached(video_path: str, cap: cv2.VideoCapture) -> list:
    """
    ffprobe を使用してキーフレームを取得。
    キャッシュ機構で重複呼び出しを回避。
    フォールバック時も同じ VideoCapture インスタンスを活用。
    """
    cmd = [
        "ffprobe",
        "-v", "error",
        "-select_streams", "v:0",
        "-show_entries", "packet=pts_time,flags",
        "-of", "csv=print_section=0",
        video_path
    ]
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=60  # タイムアウト短縮
        )
        timestamps = []
        for line in result.stdout.splitlines():
            parts = line.strip().split(",")
            if len(parts) >= 2 and "K" in parts[1].upper():
                try:
                    timestamps.append(float(parts[0]))
                except ValueError:
                    pass
        if timestamps:
            return sorted(set(timestamps))
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # フォールバック: 1秒間隔でサンプリング
    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    step = max(1, int(fps))
    return [i / fps for i in range(0, total_frames, step)]


def detect_static_scenes(
    video_path: str,
    change_level: int = 5,
    task_state = None
) -> tuple:
    """
    I フレーム（キーフレーム）を順に参照し、基準フレームとの差分が
    閾値以上になったフレームを新しいスライドとして抽出する。

    最初に読めた I フレームを 1 枚目のスライド兼基準フレームとし、
    以降は「現在フレーム」と「基準フレーム」の MAD を比較する。
    差分が閾値以上ならスライドが変わったと判定し、そのフレームを保存して
    次の比較基準に更新する。

    Parameters
    ----------
    video_path : str
        解析対象の動画ファイルパス。
    change_level : int
        変化検知の感度レベル（1〜10）。
        1 = 非常に敏感（わずかな変化も検知）、
        10 = 鈍感（大きな変化のみ検知）。

    Returns
    -------
    tuple[list, str]
        (scenes, task_temp_dir)
        scenes: [{"timestamp": float, "image_path": str}, ...]
    """
    # 変化レベルを MAD 閾値に変換
    change_level = max(1, min(10, int(change_level)))
    threshold = _CHANGE_LEVEL_TO_THRESHOLD[change_level]
    print(f"変化レベル: {change_level} → MAD 閾値: {threshold:.1f}")
    print("検出方式: 基準フレームとの差分でスライド切替を判定")

    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise RuntimeError("動画ファイルを開けませんでした。")

    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    duration = total_frames / fps if fps > 0 else 0
    if fps <= 0:
        cap.release()
        raise RuntimeError("動画の FPS を取得できませんでした。")

    print(f"動画解析開始: FPS={fps:.2f}, 総フレーム数={total_frames}, 長さ={duration:.2f}秒")

    # I フレームのタイムスタンプを取得（改善版）
    print("I フレーム（キーフレーム）のタイムスタンプを取得中...")
    keyframe_times = get_keyframe_timestamps_cached(video_path, cap)
    print(f"参照する I フレーム数: {len(keyframe_times)}")

    scenes = []
    task_id = str(uuid.uuid4())
    task_temp_dir = os.path.join(TEMP_DIR, task_id)
    os.makedirs(task_temp_dir, exist_ok=True)

    def save_scene(frame, timestamp: float) -> None:
        slide_index = len(scenes)
        img_name = f"slide_{slide_index}_{int(timestamp)}.jpg"
        img_path = os.path.join(task_temp_dir, img_name)
        cv2.imwrite(img_path, frame)
        scenes.append({
            "timestamp": timestamp,
            "image_path": img_path
        })

    base_gray_small = None
    processed_count = 0
    changed_count = 0

    print("各Iフレームを基準フレームと比較中...")
    try:
        for i, ts in enumerate(keyframe_times):
            if task_state:
                task_state.check_cancelled()
                # 30% から 80% の間で進捗を計算
                progress = 30 + int((i + 1) / len(keyframe_times) * 50)
                if (i + 1) % max(1, len(keyframe_times) // 10) == 0 or i == 0 or i == len(keyframe_times) - 1:
                    task_state.log(f"スライド画像を抽出中... ({i+1}/{len(keyframe_times)} フレーム解析)", progress)

            if (i + 1) % max(1, len(keyframe_times) // 10) == 0:
                print(f"  進捗: {i+1}/{len(keyframe_times)} ({(i+1)*100//len(keyframe_times)}%)")

            frame_idx = int(ts * fps)
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            ret, frame = cap.read()
            if not ret:
                continue

            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            gray_small = cv2.resize(gray, (160, 90))
            processed_count += 1

            if base_gray_small is None:
                save_scene(frame, ts)
                base_gray_small = gray_small
                print(f"  初期スライド抽出: {ts:.2f}秒")
                continue

            mad = float(np.mean(cv2.absdiff(gray_small, base_gray_small)))
            if mad >= threshold:
                save_scene(frame, ts)
                base_gray_small = gray_small
                changed_count += 1
                print(f"  スライド切替検出: {ts:.2f}秒 (MAD={mad:.1f})")
    finally:
        cap.release()

    print(f"フレーム処理完了: {processed_count} フレーム処理, {changed_count} 回の切替を検出")

    print(f"抽出完了: {len(scenes)} 枚のスライド画像を抽出しました")
    return scenes, task_temp_dir


def create_presentation(
    title: str,
    scenes: list,
    transcript: list,
    output_pptx_path: str,
    url: str | None = None,
    ai_summary_enabled: bool = False,
    task_state = None
):
    """
    抽出した画像と文字起こしテキストをマッピングし、PowerPointプレゼンテーションを生成する。
    transcript が空リストの場合は画像のみのスライドを生成する。
    
    ai_summary_enabled が True の場合のみ、スライドごとに字幕をまとめて
    5行の要点 + メインメッセージを抽出し、スライドに追加する。
    """
    prs = Presentation()
    # スライドサイズを 16:9 ワイドスクリーンに設定
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 白紙スライドのレイアウト (blank layout is index 6)
    blank_slide_layout = prs.slide_layouts[6]

    summarizer = None
    if ai_summary_enabled:
        try:
            summarizer = GeminiSummarizer()
            print("✓ Gemini API が有効です。スライド要約を生成します。")
        except Exception as e:
            ai_summary_enabled = False
            print(f"ℹ Gemini 要約を無効化しました。({e})")

    # --- 1. タイトルスライドの作成 ---
    slide = prs.slides.add_slide(blank_slide_layout)

    # タイトル用テキストボックス
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(33, 37, 41)

    p2 = tf.add_paragraph()
    p2.text = "YouTube動画から自動生成されたスライド資料"
    p2.font.size = Pt(20)
    p2.font.name = "Arial"

    # URL を追加（存在する場合）
    if url:
        p3 = tf.add_paragraph()
        p3.text = f"URL: {url}"
        p3.font.size = Pt(14)
        p3.font.name = "Arial"
        p3.font.color.rgb = RGBColor(0, 102, 204)  # 目立つ色
    p2.font.color.rgb = RGBColor(108, 117, 125)
    p2.space_before = Pt(20)

    # --- 2. コンテンツスライドの作成 ---
    has_transcript = bool(transcript)

    for i, scene in enumerate(scenes):
        if task_state:
            task_state.check_cancelled()
            # 80% から 95% の間で進捗を計算
            progress = 80 + int((i + 1) / len(scenes) * 15)
            task_state.log(f"PowerPointスライドを作成中... ({i+1}/{len(scenes)}枚目)", progress)

        current_time = scene["timestamp"]
        next_time = scenes[i + 1]["timestamp"] if i + 1 < len(scenes) else float('inf')

        slide = prs.slides.add_slide(blank_slide_layout)

        if has_transcript:
            # 左側: 画像（幅 5.8インチ）
            img_left = Inches(0.8)
            img_top = Inches(1.5)
            img_width = Inches(5.8)
            slide.shapes.add_picture(scene["image_path"], img_left, img_top, width=img_width)

            # このスライドの時間範囲に該当する字幕を結合
            slide_text = get_scene_text(transcript, current_time, next_time)
            if not slide_text:
                slide_text = "(この区間の文字起こしデータはありません)"

            # 右側: テキストボックス
            text_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.5))
            tf = text_box.text_frame
            tf.word_wrap = True

            p_time = tf.paragraphs[0]
            p_time.text = f"【シーン開始: {format_timestamp(current_time)}】"
            p_time.font.size = Pt(14)
            p_time.font.bold = True
            p_time.font.color.rgb = RGBColor(0, 123, 255)
            p_time.space_after = Pt(10)

            # === Gemini で要約を生成 ===
            if ai_summary_enabled and summarizer and slide_text and "(この区間の文字起こしデータはありません)" not in slide_text:
                try:
                    summary_result = summarizer.summarize_slide_content(slide_text)

                    # 【キーポイント】セクション
                    if summary_result.get("key_points"):
                        p_key_title = tf.add_paragraph()
                        p_key_title.text = "【要点】"
                        p_key_title.font.size = Pt(12)
                        p_key_title.font.bold = True
                        p_key_title.font.color.rgb = RGBColor(220, 53, 69)
                        p_key_title.space_before = Pt(8)
                        p_key_title.space_after = Pt(4)

                        for key_point in summary_result["key_points"]:
                            p_kp = tf.add_paragraph()
                            p_kp.text = f"• {format_slide_text(key_point)}"
                            p_kp.font.size = Pt(11)
                            p_kp.font.name = "Arial"
                            p_kp.font.color.rgb = RGBColor(50, 50, 50)
                            p_kp.level = 0
                            p_kp.space_after = Pt(3)

                    # 【メインメッセージ】セクション
                    if summary_result.get("main_message"):
                        p_main_title = tf.add_paragraph()
                        p_main_title.text = "💡 最も言いたいこと"
                        p_main_title.font.size = Pt(12)
                        p_main_title.font.bold = True
                        p_main_title.font.color.rgb = RGBColor(0, 102, 204)
                        p_main_title.space_before = Pt(8)
                        p_main_title.space_after = Pt(4)

                        p_main = tf.add_paragraph()
                        p_main.text = format_slide_text(summary_result["main_message"])
                        p_main.font.size = Pt(11)
                        p_main.font.name = "Arial"
                        p_main.font.color.rgb = RGBColor(20, 20, 100)
                        p_main.font.italic = True
                        p_main.line_spacing = 1.2

                except Exception as e:
                    print(f"スライド {i+1} の要約生成に失敗: {e}")
                    # フォールバック: 元の字幕を表示
                    p_content = tf.add_paragraph()
                    p_content.text = slide_text
                    p_content.font.size = Pt(11)
                    p_content.font.name = "Arial"
                    p_content.font.color.rgb = RGBColor(50, 50, 50)
                    p_content.line_spacing = 1.3
            else:
                # Geminiがない場合は元の字幕を表示
                p_content = tf.add_paragraph()
                p_content.text = slide_text
                p_content.font.size = Pt(11)
                p_content.font.name = "Arial"
                p_content.font.color.rgb = RGBColor(50, 50, 50)
                p_content.line_spacing = 1.3

        else:
            # 字幕なし: 画像をスライド全体に広げて配置
            img_left = Inches(0.8)
            img_top = Inches(0.8)
            img_width = Inches(11.733)
            slide.shapes.add_picture(scene["image_path"], img_left, img_top, width=img_width)

            # タイムスタンプのみ表示
            ts_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(4.0), Inches(0.5))
            tf = ts_box.text_frame
            p = tf.paragraphs[0]
            p.text = format_timestamp(current_time)
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(150, 150, 150)

    prs.save(output_pptx_path)
    print(f"PowerPointファイルを保存しました: {output_pptx_path}")


def process_youtube_to_presentation(
    url: str,
    change_level: int = 5,
    ai_summary_enabled: bool = False,
    save_format: str = "pptx",
    task_state = None
) -> dict:
    """
    YouTube URL からプレゼンテーションを生成する一連の処理を実行する。
    字幕が取得できない場合は画像のみのスライドを生成する。

    Parameters
    ----------
    url : str
        YouTube 動画の URL。
    change_level : int
        変化検知の感度（1〜10）。1が最も敏感、10が最も鈍感。
    ai_summary_enabled : bool
        True のときだけ Gemini 要約を有効化する。
    """
    video_id = extract_video_id(url)
    
    # 安全対策として、抽出した動画IDから正規のYouTube URLを再構築する
    url = f"https://www.youtube.com/watch?v={video_id}"

    # 1. 字幕の取得（失敗しても処理続行）
    if task_state:
        task_state.log("YouTube動画の字幕を取得中...", 5)
    print("字幕の取得中...")
    transcript = get_transcript(url, video_id)
    if transcript:
        print(f"字幕を取得しました（{len(transcript)} エントリ）。")
        if task_state:
            task_state.log(f"字幕を取得しました（{len(transcript)} エントリ）。", 10)
    else:
        print("字幕なし。画像のみのスライドを生成します。")
        if task_state:
            task_state.log("字幕が見つかりませんでした。画像のみのスライドを生成します。", 10)

    # 動画IDに対応した固定の保存パスを設定（再ダウンロード防止）
    temp_video_path = os.path.join(TEMP_DIR, f"video_{video_id}.mp4")
    task_temp_dir = None

    try:
        # 2. 動画のダウンロード（存在する場合はスキップ）
        if task_state:
            task_state.log("動画のダウンロードを開始しました...", 15)
        print("動画のダウンロード中...")
        title = download_video(url, temp_video_path)
        if task_state:
            task_state.log(f"動画のダウンロードが完了しました: {title}", 30)

        # 3. I フレームを参照したスライド切替の検知と画像抽出
        if task_state:
            task_state.log("動画の解析（Iフレーム抽出）を開始しました...", 30)
        print("スライド切替の検知中（I フレームのみ参照）...")
        scenes, task_temp_dir = detect_static_scenes(
            temp_video_path,
            change_level=change_level,
            task_state=task_state
        )

        if not scenes:
            raise ValueError(
                "スライド画像が検出されませんでした。"
                "変化レベルを下げて、より小さな変化も検出するようにしてください。"
            )

        # スライド枚数が100枚を超えたら、確認ダイアログを表示するため一時停止する
        if len(scenes) > 100:
            if task_state:
                task_state.status = "waiting_confirm"
                task_state.log(f"スライド数が100枚を超えました（現在 {len(scenes)} 枚）。一時停止中...", task_state.progress)
                task_state.confirm_event.clear()
                task_state.confirm_event.wait()
                if task_state.confirm_response == "abort":
                    raise TaskCancelledException("スライド数が100枚を超えたため、ユーザーによって中止されました。")
                else:
                    task_state.status = "processing"
                    task_state.log("処理を継続します。", task_state.progress)

        # 文字起こしテキストが日本語でない場合の翻訳確認フロー
        full_text = " ".join([e["text"] for e in transcript]) if transcript else ""
        if transcript and not is_japanese(full_text):
            if task_state:
                task_state.status = "waiting_translation"
                # フロントエンドに字幕データを渡すため、一時的に result に格納する
                task_state.result = {
                    "transcript": transcript
                }
                task_state.log("文字起こしテキストの翻訳を確認しています...", task_state.progress)
                task_state.confirm_event.clear()
                task_state.confirm_event.wait()

                # 中止された場合は処理を抜ける
                if task_state.confirm_response == "abort":
                    raise TaskCancelledException("翻訳確認画面でユーザーによって中止されました。")
                elif task_state.confirm_response == "use_translation":
                    translated_texts = getattr(task_state, "translated_texts", None)
                    if translated_texts and len(translated_texts) == len(transcript):
                        for idx, text in enumerate(translated_texts):
                            # 各字幕テキストを翻訳されたものに置き換え、日本語整形処理を適用する
                            transcript[idx]["text"] = format_slide_text(text)
                        task_state.log("表示された翻訳テキストを整形して適用しました。", task_state.progress)
                    else:
                        task_state.log("警告: 翻訳されたテキストの数が一致しないため、置換されませんでした。", task_state.progress)
                else:
                    task_state.log("元のテキスト（英語のまま）で処理を継続します。", task_state.progress)

                # ステータスを processing に戻し、一時的な result をクリア
                task_state.status = "processing"
                task_state.result = None

        # 4. 出力ファイルの生成
        # ファイル名は動画タイトルをベースに生成（特殊文字を除去）
        safe_title = sanitize_filename(title)
        markdown_asset_dirname = None
        markdown_asset_filenames = []
        if save_format == "markdown":
            print("Markdownファイルと画像フォルダの生成中...")
            if task_state:
                task_state.log("Markdown資料の生成を開始しました...", 80)
            output_filename, markdown_asset_dirname, markdown_asset_filenames = create_markdown_package(
                title=title,
                scenes=scenes,
                transcript=transcript,
                task_temp_dir=task_temp_dir,
                safe_title=safe_title,
                url=url,
                ai_summary_enabled=ai_summary_enabled,
                task_state=task_state,
            )
            download_label = "Markdown 一式を保存"
        else:
            output_filename = f"{safe_title}.pptx"
            output_pptx_path = os.path.join(task_temp_dir, output_filename)

            counter = 1
            base_name = output_filename.replace(".pptx", "")
            while os.path.exists(output_pptx_path):
                output_filename = f"{base_name}_{counter}.pptx"
                output_pptx_path = os.path.join(task_temp_dir, output_filename)
                counter += 1

            print("PowerPointの生成中...")
            if task_state:
                task_state.log("PowerPointプレゼンテーションの生成を開始しました...", 80)
            create_presentation(
                title,
                scenes,
                transcript,
                output_pptx_path,
                url=url,
                ai_summary_enabled=ai_summary_enabled,
                task_state=task_state,
            )
            download_label = "PowerPoint をダウンロード"

        # フロントエンドに返す結果データを整形
        if task_state:
            task_state.log("結果データを生成中...", 95)
        formatted_scenes = []
        for i, scene in enumerate(scenes):
            current_time = scene["timestamp"]
            next_time = scenes[i + 1]["timestamp"] if i + 1 < len(scenes) else float('inf')

            # テキストの再抽出（プレビュー用）
            slide_text = get_scene_text(transcript, current_time, next_time)

            formatted_scenes.append({
                "id": i,
                "timestamp": current_time,
                "image_name": os.path.basename(scene["image_path"]),
                "text": slide_text or "(字幕なし)"
            })

        result = {
            "success": True,
            "task_id": os.path.basename(task_temp_dir),
            "title": title,
            "download_filename": output_filename,
            "download_label": download_label,
            "save_format": save_format,
            "asset_dirname": markdown_asset_dirname,
            "asset_filenames": markdown_asset_filenames,
            "scenes": formatted_scenes,
            "has_transcript": bool(transcript),
        }

        if task_state:
            task_state.log("成果物の生成が完了しました！", 100)

        return result

    except Exception as e:
        if task_temp_dir and os.path.exists(task_temp_dir):
            try:
                shutil.rmtree(task_temp_dir)
                print(f"エラー発生により一時ディレクトリをクリーンアップしました: {task_temp_dir}")
            except Exception as clean_err:
                print(f"一時ディレクトリのクリーンアップに失敗: {clean_err}")
        raise e
